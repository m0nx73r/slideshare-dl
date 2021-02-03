#!/usr/bin/env python3

import requests
from bs4 import BeautifulSoup
import os
import shutil
from PIL import Image
from pptx import Presentation   
from pptx.util import Inches
import sys
from termcolor import colored

banner = """ 
  _________.__   .__     .___       ________                          
 /   _____/|  |  |__|  __| _/ ____  \______ \    ____ __  _  __ ____            
 \_____  \ |  |  |  | / __ |_/ __ \  |    |  \  /  _ \\ \/ \/ //    \       Author:  m0nxt3r
 /        \|  |__|  |/ /_/ |\  ___/  |    `   \(  <_> )\     /|   |  \     Version: 0.0.1
/_______  /|____/|__|\____ | \___  >/_______  / \____/  \/\_/ |___|  /     Telegram: https://t.me/m0nxt3r
        \/                \/     \/         \/                     \/        
"""


GREEN = '\033[91m'
print(GREEN + banner)

def help():
    helptext = """USAGE : python3 slidedown.py <url>"""
    print(helptext)
    print('\033[0m')    
    sys.exit(0)

if len(sys.argv) == 2:
    try:
        url = sys.argv[1]
        url = url.split("?")[0]
        print(colored("[+] PROCESSING URL: ", "white", attrs=["bold"]), colored(url, "cyan"))
    except IndexError:
        help()
else:
    help()

dirname = "." + url.split("/")[-1]

def scrape(url):
    r = requests.get(url)
    soup = BeautifulSoup(r.content, "lxml")
    downloader(soup)

def downloader(soup):
    slide = soup.findAll("img", {'class': 'slide_image'})
    global num
    num = len(slide)
    try:
        os.mkdir(dirname)
    except FileExistsError:
        shutil.rmtree(dirname)
        os.mkdir(dirname)
    for i in range(0, num):
        r = requests.get(slide[i]["data-full"])
        print("[+] Downloading slide " + str(i+1), end="\r")
        with open(dirname + "/" + str(i) + ".jpg", "wb") as f:
            f.write(r.content)
            f.close()
    makeprs()

def _add_image(slide, image_url):
    placeholder = slide.shapes
 
    # Calculate the image size of the image
    im = Image.open(image_url)
    width, height = im.size
 
    # Make sure the placeholder doesn't zoom in
    placeholder.height = height
    placeholder.width = width   
    
    # Insert the picture
    placeholder = placeholder.add_picture(image_url, left=Inches(0) , top=Inches(0), width=Inches(10.6) , height=Inches(8))
 
    # Calculate ratios and compare
    image_ratio = width / height
    placeholder_ratio = placeholder.width / placeholder.height
    ratio_difference = placeholder_ratio - image_ratio
 
    # Placeholder width too wide:
    if ratio_difference > 0:
        difference_on_each_side = ratio_difference / 2
        placeholder.crop_left = -difference_on_each_side
        placeholder.crop_right = -difference_on_each_side
    # Placeholder height too high
    else:
        difference_on_each_side = -ratio_difference / 2
        placeholder.crop_bottom = -difference_on_each_side
        placeholder.crop_top = -difference_on_each_side

def makeprs():
    print("[+] Processing your presentation...")          
    prs = Presentation()
    prs.slide_height = Inches(8)
    prs.slide_width = Inches(10.666667) 
    for _ in range(num):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

    for i in range(num):
        _add_image(prs.slides[i],dirname + "/" + str(i) + ".jpg")

    prs.save(dirname.replace(".", "") + ".pptx")
    shutil.rmtree(dirname)
    print("[+] PPT saved as " + dirname + ".pptx in " + os.getcwd() + ".")

def start(url):
    scrape(url)

if __name__ == "__main__":
    start(url)

